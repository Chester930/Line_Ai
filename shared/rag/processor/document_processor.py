import os
import logging
from typing import Dict, Any, Optional, List
from datetime import datetime
import hashlib
import mimetypes
from pathlib import Path
import PyPDF2
import docx
import pandas as pd
import io
from pptx import Presentation
import win32com.client
import pythoncom
import tempfile
import shutil

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """文件處理器基礎類"""
    
    def __init__(self):
        self.supported_types = {
            'text/plain': self._process_text,
            'application/pdf': self._process_pdf,
            'application/msword': self._process_doc,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx,
            'application/vnd.ms-excel': self._process_excel,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_xlsx,
            'application/vnd.ms-powerpoint': self._process_ppt,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx
        }
        self.max_file_size = 200 * 1024 * 1024  # 200MB
        
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """處理文件的主要方法"""
        try:
            # 檢查文件是否存在
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            # 檢查文件大小
            file_size = os.path.getsize(file_path)
            if file_size > self.max_file_size:
                raise ValueError(f"文件超過大小限制 ({self.max_file_size/1024/1024}MB)")
            
            # 獲取文件類型
            mime_type = mimetypes.guess_type(file_path)[0]
            if not mime_type or mime_type not in self.supported_types:
                raise ValueError(f"不支持的文件類型: {mime_type}")
            
            # 處理文件
            processor = self.supported_types[mime_type]
            content = processor(file_path)
            
            # 生成文件信息
            file_info = {
                'file_path': file_path,
                'file_name': os.path.basename(file_path),
                'file_type': mime_type,
                'file_size': file_size,
                'content': content,
                'content_hash': self._calculate_hash(content),
                'processed_at': datetime.utcnow().isoformat(),
                'metadata': {
                    'original_file': {
                        'name': os.path.basename(file_path),
                        'extension': Path(file_path).suffix,
                        'created_at': datetime.fromtimestamp(os.path.getctime(file_path)).isoformat(),
                        'modified_at': datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat()
                    }
                }
            }
            
            return file_info
            
        except Exception as e:
            logger.error(f"處理文件時發生錯誤: {str(e)}")
            raise
    
    def _calculate_hash(self, content: str) -> str:
        """計算內容的 SHA-256 hash"""
        return hashlib.sha256(content.encode('utf-8')).hexdigest()
    
    def _process_text(self, file_path: str) -> str:
        """處理純文本文件"""
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _process_pdf(self, file_path: str) -> str:
        """處理 PDF 文件"""
        try:
            text_content = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # 提取文檔信息
                metadata = pdf_reader.metadata
                if metadata:
                    text_content.append(f"文檔信息:")
                    for key, value in metadata.items():
                        if value:
                            text_content.append(f"{key}: {value}")
                    text_content.append("\n")
                
                # 提取每頁內容
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text_content.append(f"\n=== 第 {page_num} 頁 ===\n")
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                
            return "\n".join(text_content)
            
        except Exception as e:
            logger.error(f"處理 PDF 文件時發生錯誤: {str(e)}")
            raise
    
    def _process_doc(self, file_path: str) -> str:
        """處理 DOC 文件"""
        # TODO: 實現 DOC 處理邏輯
        raise NotImplementedError("DOC 處理功能尚未實現")
    
    def _process_docx(self, file_path: str) -> str:
        """處理 DOCX 文件"""
        try:
            doc = docx.Document(file_path)
            content_parts = []
            
            # 提取標題
            if doc.core_properties.title:
                content_parts.append(f"標題: {doc.core_properties.title}\n")
            
            # 提取段落
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)
            
            # 提取表格
            for table in doc.tables:
                content_parts.append("\n表格內容:")
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    content_parts.append(" | ".join(row_text))
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"處理 DOCX 文件時發生錯誤: {str(e)}")
            raise
    
    def _process_excel(self, file_path: str) -> str:
        """處理 XLS 文件"""
        # XLS 文件使用相同的處理邏輯
        return self._process_xlsx(file_path)
    
    def _process_xlsx(self, file_path: str) -> str:
        """處理 XLSX 文件"""
        try:
            content_parts = []
            
            # 讀取所有工作表
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                content_parts.append(f"\n=== 工作表: {sheet_name} ===\n")
                
                # 獲取列名
                content_parts.append("列名: " + " | ".join(df.columns.astype(str)))
                
                # 轉換數據為文本描述
                for index, row in df.iterrows():
                    row_text = " | ".join(row.astype(str))
                    content_parts.append(row_text)
                
                # 添加基本統計信息
                content_parts.append(f"\n基本統計:")
                content_parts.append(f"行數: {len(df)}")
                content_parts.append(f"列數: {len(df.columns)}")
                
                # 對數值列進行統計
                numeric_columns = df.select_dtypes(include=['int64', 'float64']).columns
                if not numeric_columns.empty:
                    content_parts.append("\n數值列統計:")
                    stats = df[numeric_columns].describe().to_string()
                    content_parts.append(stats)
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"處理 XLSX 文件時發生錯誤: {str(e)}")
            raise
    
    def _process_ppt(self, file_path: str) -> str:
        """處理 PPT 文件"""
        try:
            # 使用 win32com 將 PPT 轉換為 PPTX
            pythoncom.CoInitialize()
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            
            # 創建臨時文件
            temp_dir = tempfile.mkdtemp()
            temp_pptx = os.path.join(temp_dir, "temp.pptx")
            
            try:
                # 打開並保存為 PPTX
                presentation = powerpoint.Presentations.Open(file_path)
                presentation.SaveAs(temp_pptx, 24)  # 24 是 PPTX 格式的常量
                presentation.Close()
                
                # 處理轉換後的 PPTX
                content = self._process_pptx(temp_pptx)
                return content
                
            finally:
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                # 清理臨時文件
                shutil.rmtree(temp_dir, ignore_errors=True)
                
        except Exception as e:
            logger.error(f"處理 PPT 文件時發生錯誤: {str(e)}")
            raise
    
    def _process_pptx(self, file_path: str) -> str:
        """處理 PPTX 文件"""
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            # 提取文檔屬性
            if prs.core_properties.title:
                content_parts.append(f"標題: {prs.core_properties.title}")
            if prs.core_properties.subject:
                content_parts.append(f"主題: {prs.core_properties.subject}")
            if prs.core_properties.author:
                content_parts.append(f"作者: {prs.core_properties.author}")
            content_parts.append("")
            
            # 處理每一張幻燈片
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n=== 幻燈片 {slide_num} ===\n")
                
                # 提取幻燈片標題
                if slide.shapes.title:
                    content_parts.append(f"標題: {slide.shapes.title.text}")
                
                # 提取所有文本框和表格
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 文本內容
                        content_parts.append(shape.text.strip())
                    
                    elif shape.has_table:
                        # 表格內容
                        content_parts.append("\n表格內容:")
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            content_parts.append(" | ".join(row_text))
                
                # 添加註釋（如果有）
                if hasattr(slide, "notes_slide") and slide.notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        content_parts.append(f"\n註釋:\n{notes}")
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"處理 PPTX 文件時發生錯誤: {str(e)}")
            raise
